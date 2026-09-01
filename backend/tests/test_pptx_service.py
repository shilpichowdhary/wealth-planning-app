"""Deck rendering: mandatory slides + brand-font embedding (feedback #6)."""
import zipfile

import pytest
from pptx import Presentation

from backend.services.pptx_service import build_pptx
from backend.services.font_embed import embed_fonts, SKILL_FONT_DIR, _resolve_fonts


def _minimal_spec():
    return {
        "client_name": "Test Client",
        "slides": [
            {"layout": "cover", "client_name": "Test Client", "title": "Advisory review."},
            {"layout": "07", "eyebrow": "Recommendation", "title": "A single clear statement.",
             "lede": "Supporting paragraph that explains the recommendation in a couple of sentences."},
        ],
    }


def test_disclaimer_and_offices_auto_appended(tmp_path):
    out = build_pptx(_minimal_spec(), tmp_path / "deck.pptx")
    prs = Presentation(str(out))
    # 2 authored slides + auto Disclaimer + auto Offices = 4
    assert len(prs.slides) == 4

    all_text = "\n".join(
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Disclaimer" in all_text
    assert "Offices" in all_text
    # Regulatory footer must be present verbatim.
    assert "MAS" in all_text and "DFSA" in all_text and "SEBI" in all_text and "FCA" in all_text
    # Accredited Investors / Professional Clients scoping line.
    assert "Accredited Investors" in all_text


def test_disclaimer_not_duplicated_when_present(tmp_path):
    spec = _minimal_spec()
    spec["slides"].append({"layout": "04"})  # author already included disclaimer
    out = build_pptx(spec, tmp_path / "deck2.pptx")
    prs = Presentation(str(out))
    texts = [
        shape.text_frame.text
        for slide in prs.slides for shape in slide.shapes
        if shape.has_text_frame
    ]
    # "Disclaimer" heading appears on exactly one slide, not two.
    disclaimer_slides = sum(
        1 for slide in prs.slides
        if any(s.has_text_frame and "intended recipient" in s.text_frame.text
               for s in slide.shapes)
    )
    assert disclaimer_slides == 1


@pytest.mark.skipif(not _resolve_fonts(), reason="brand font files not present")
def test_fonts_are_embedded(tmp_path):
    out = build_pptx(_minimal_spec(), tmp_path / "deck.pptx")
    with zipfile.ZipFile(out, "r") as z:
        names = z.namelist()
        assert any(n.startswith("ppt/fonts/") and n.endswith(".fntdata") for n in names), \
            "no embedded font parts found"
        presentation = z.read("ppt/presentation.xml").decode("utf-8")
        content_types = z.read("[Content_Types].xml").decode("utf-8")
        rels = z.read("ppt/_rels/presentation.xml.rels").decode("utf-8")
    assert 'embedTrueTypeFonts="1"' in presentation
    assert "<p:embeddedFontLst>" in presentation
    assert 'typeface="Public Sans"' in presentation
    assert 'typeface="Frank Ruhl Libre"' in presentation
    assert 'Extension="fntdata"' in content_types
    assert "/font" in rels  # font relationship type present


def test_embed_fonts_is_idempotent(tmp_path):
    out = build_pptx(_minimal_spec(), tmp_path / "deck.pptx")
    # build_pptx already embedded once; a second explicit pass must not
    # duplicate the content-type default or corrupt the package.
    embed_fonts(out)
    prs = Presentation(str(out))  # still opens
    with zipfile.ZipFile(out, "r") as z:
        content_types = z.read("[Content_Types].xml").decode("utf-8")
    assert content_types.count('Extension="fntdata"') == 1
