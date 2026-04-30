"""LC-branded PowerPoint deck rendering.

Takes the JSON slide spec emitted by deck_curator and produces an editable
.pptx file using python-pptx. Each spec slide carries a `layout` field
matching one of the 15 LC layouts in SKILL.md §3 — this module dispatches
to one builder per layout, plus a custom `structure` builder that embeds
the React-Flow-style diagram as a high-res PNG.

Brand contract (from SKILL.md §2):
  • Frank Ruhl Libre Light (300) for titles/labels/quotes
  • Public Sans 400/500 for body
  • Red #E50025 ONLY as 3-4px vertical rule, region underline, or italic accent
  • White backgrounds; black section dividers
  • No gradients, shadows, rounded containers, emoji, decorative icons

Conversion: 1 px @ 96 DPI = 9525 EMU; 1 pt = 12700 EMU. Slide is 1920×1080
px → 18,288,000 × 10,287,000 EMU = 20" × 11.25".
"""
from __future__ import annotations

import logging
from datetime import datetime, timezone
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

logger = logging.getLogger(__name__)

# ── Brand constants ───────────────────────────────────────────────────

LC_RED = RGBColor(0xE5, 0x00, 0x25)
LC_BLACK = RGBColor(0x00, 0x00, 0x00)
LC_BLACK_80 = RGBColor(0x33, 0x33, 0x33)
LC_BLACK_60 = RGBColor(0x66, 0x66, 0x66)
LC_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LC_SMOKE = RGBColor(0xFA, 0xFA, 0xFA)
LC_INK_400 = RGBColor(0xCD, 0xCD, 0xCD)

FONT_SERIF = "Frank Ruhl Libre"
FONT_SANS = "Public Sans"

# Slide canvas — 1920×1080 @ 96 DPI in EMU
SLIDE_W_EMU = Emu(1920 * 9525)
SLIDE_H_EMU = Emu(1080 * 9525)

# Safe zone — 120px from each edge per SKILL.md §1
SAFE_PAD = 120


def _px(n: float) -> Emu:
    """Pixels → EMU (96 DPI)."""
    return Emu(int(n * 9525))


# ── Shape helpers ─────────────────────────────────────────────────────


def _add_textbox(
    slide,
    *,
    x_px: float,
    y_px: float,
    w_px: float,
    h_px: float,
    text: str,
    font: str,
    size_pt: float,
    color: RGBColor = LC_BLACK,
    bold: bool = False,
    italic: bool = False,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    letter_spacing_em: float | None = None,
    upper: bool = False,
):
    """Add a text box with the LC brand font/size baked in.

    `letter_spacing_em` in fractional em (e.g. 0.18) — converted to font.spacing
    (1/100 pt). `upper` capitalises the text (for eyebrows). We never set
    bold=True on a Frank Ruhl Libre run because Light(300) is the only weight
    we use; bold on body Public Sans is allowed at 600+ via weight.
    """
    box = slide.shapes.add_textbox(_px(x_px), _px(y_px), _px(w_px), _px(h_px))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if upper else text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if letter_spacing_em:
        # python-pptx exposes character spacing via XML directly. Spacing is
        # in 100ths of a point; convert from em (relative to font size).
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(letter_spacing_em * size_pt * 100)))
    return box


def _add_rect(slide, *, x_px: float, y_px: float, w_px: float, h_px: float, fill: RGBColor):
    """Solid-fill rectangle, no border. Used for the 3-4px red rule device."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _px(x_px), _px(y_px), _px(w_px), _px(h_px)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    rect.line.fill.background()
    return rect


def _add_line(slide, *, x_px: float, y_px: float, w_px: float, color: RGBColor, weight_pt: float = 2):
    """Horizontal line — used for region underlines and section dividers."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, _px(x_px), _px(y_px), _px(w_px), Emu(int(weight_pt * 12700))
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def _new_slide(prs: Presentation, *, bg: RGBColor = LC_WHITE):
    """Add a blank slide with the requested background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank
    bg_shape = slide.background
    fill = bg_shape.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def _chapter_tab(slide, *, num: str, label: str):
    """Top-of-slide chapter tab — small uppercase text + serif red number.

    Mirrors the .chapter-tab pattern in ppt-template.css: positioned at
    top:56px, left:120px, used on every content slide so the deck reads as
    a coherent navigated document.
    """
    _add_textbox(
        slide,
        x_px=SAFE_PAD,
        y_px=56,
        w_px=200,
        h_px=24,
        text=num,
        font=FONT_SERIF,
        size_pt=14,
        color=LC_RED,
    )
    _add_textbox(
        slide,
        x_px=SAFE_PAD + 38,
        y_px=58,
        w_px=1500,
        h_px=20,
        text=label,
        font=FONT_SANS,
        size_pt=11,
        color=LC_BLACK_80,
        upper=True,
        letter_spacing_em=0.14,
    )


def _page_no(slide, n: int):
    """Page number, footer-right, serif tabular-nums."""
    _add_textbox(
        slide,
        x_px=1920 - SAFE_PAD - 60,
        y_px=1080 - 76,
        w_px=60,
        h_px=24,
        text=f"{n:02d}",
        font=FONT_SERIF,
        size_pt=14,
        color=LC_BLACK,
        align=PP_ALIGN.RIGHT,
    )


# ── Layout 02: Cover (typography-only, white bg) ──────────────────────


def _build_cover(prs: Presentation, slide_spec: dict, page_no: int = 1):
    """Layout 02 — cover-alt. White bg, big serif title, red rule, prepared-for
    block. We use this for case decks (no photo available). page_no ignored
    on the cover; signature matches other builders for uniform dispatch.

    Spec fields:
      title    — big serif statement (defaults to 'Wealth Planning Advisory')
      subtitle — optional smaller serif line below the title (the case-specific
                 framing, e.g. 'Singapore trust and Section 13U capital flow')
    """
    slide = _new_slide(prs)
    client = slide_spec.get("client_name", "Confidential Client")
    title = slide_spec.get("title") or "Wealth Planning Advisory"
    subtitle = slide_spec.get("subtitle") or ""

    _add_textbox(
        slide,
        x_px=SAFE_PAD,
        y_px=SAFE_PAD,
        w_px=1500,
        h_px=24,
        text="Lighthouse Canton  ·  Wealth Planning",
        font=FONT_SANS,
        size_pt=11,
        color=LC_RED,
        upper=True,
        letter_spacing_em=0.18,
    )

    # Big title — single statement ending in a period, per LC voice rules.
    # Long titles (curator sometimes produces >40 chars) wrap to 2-3 lines at
    # 72pt; we shrink the size adaptively and reserve vertical space for the
    # full wrapped title so it never overflows into the subtitle below.
    title_text = title if title.rstrip().endswith(".") else f"{title}."
    if len(title_text) > 60:
        title_size = 56
    elif len(title_text) > 40:
        title_size = 64
    else:
        title_size = 72
    title_h = 320  # accommodates up to ~3 wrapped lines at 56-72pt
    _add_textbox(
        slide,
        x_px=SAFE_PAD,
        y_px=300,
        w_px=1600,
        h_px=title_h,
        text=title_text,
        font=FONT_SERIF,
        size_pt=title_size,
        color=LC_BLACK,
    )
    if subtitle:
        _add_textbox(
            slide, x_px=SAFE_PAD, y_px=300 + title_h + 20, w_px=1500, h_px=70,
            text=subtitle, font=FONT_SERIF, size_pt=26, color=LC_BLACK_80,
            italic=True,
        )

    # 4px red vertical rule
    _add_rect(slide, x_px=SAFE_PAD, y_px=720, w_px=4, h_px=120, fill=LC_RED)
    _add_textbox(
        slide,
        x_px=SAFE_PAD + 36,
        y_px=720,
        w_px=400,
        h_px=20,
        text="Prepared for",
        font=FONT_SANS,
        size_pt=11,
        color=LC_BLACK_80,
        upper=True,
        letter_spacing_em=0.18,
    )
    _add_textbox(
        slide,
        x_px=SAFE_PAD + 36,
        y_px=752,
        w_px=1300,
        h_px=120,
        text=client,
        font=FONT_SERIF,
        size_pt=42,
        color=LC_BLACK,
    )

    # Confidential / date stamp, footer-left
    today = datetime.now(timezone.utc).strftime("%d %B %Y")
    _add_textbox(
        slide,
        x_px=SAFE_PAD,
        y_px=1080 - 80,
        w_px=600,
        h_px=20,
        text=f"Confidential  ·  {today}",
        font=FONT_SANS,
        size_pt=11,
        color=LC_BLACK_80,
        upper=True,
        letter_spacing_em=0.14,
    )


# ── Layout 06: Two-column pillars — the hallmark LC layout ────────────


def _build_pillars(prs: Presentation, slide_spec: dict, page_no: int):
    """Layout 06 — eyebrow + serif title + N rows of (label / red-rule / body).

    The signature LC layout. spec.rows = [{label, body, eyebrow?}, ...]
    """
    slide = _new_slide(prs)
    _chapter_tab(
        slide,
        num=f"{page_no:02d}",
        label=slide_spec.get("chapter") or "Advisory analysis",
    )

    eyebrow = slide_spec.get("eyebrow") or "Findings"
    title = slide_spec.get("title") or ""
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=140, w_px=1500, h_px=24, text=eyebrow,
        font=FONT_SANS, size_pt=11, color=LC_RED, upper=True, letter_spacing_em=0.18,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=176, w_px=1600, h_px=120,
        text=title, font=FONT_SERIF, size_pt=52, color=LC_BLACK,
    )

    rows = slide_spec.get("rows") or []
    # Lay rows out vertically starting at y=340. Each row:
    #   serif label (240w) | 3px red rule | sans body (1100w)
    # Row height proportional to body length, capped to fit 4-5 rows on slide.
    row_y = 340
    row_pad = 28
    label_w = 280
    rule_x = SAFE_PAD + label_w + 24
    body_x = rule_x + 24
    body_w = 1920 - SAFE_PAD - body_x
    available_h = 1080 - row_y - 140
    n = max(1, len(rows))
    row_h = min(160, max(80, (available_h - (n - 1) * row_pad) // n))

    for row in rows:
        label = row.get("label") or row.get("title") or ""
        body = row.get("body") or ""
        # Serif label
        _add_textbox(
            slide, x_px=SAFE_PAD, y_px=row_y, w_px=label_w, h_px=row_h,
            text=label, font=FONT_SERIF, size_pt=24, color=LC_BLACK,
        )
        # 3px red rule
        _add_rect(slide, x_px=rule_x, y_px=row_y + 8, w_px=3, h_px=row_h - 16, fill=LC_RED)
        # Sans body
        _add_textbox(
            slide, x_px=body_x, y_px=row_y, w_px=body_w, h_px=row_h,
            text=body, font=FONT_SANS, size_pt=18, color=LC_BLACK_80,
        )
        row_y += row_h + row_pad

    _page_no(slide, page_no)


# ── Layout 07: Title + lede + red vertical rule ───────────────────────


def _build_title_lede(prs: Presentation, slide_spec: dict, page_no: int):
    """Layout 07 — one big serif statement + a supporting paragraph below,
    separated by a 4px red vertical rule. Use for "the call to action" or a
    single declarative claim."""
    slide = _new_slide(prs)
    _chapter_tab(
        slide, num=f"{page_no:02d}",
        label=slide_spec.get("chapter") or "Advisory analysis",
    )

    eyebrow = slide_spec.get("eyebrow") or ""
    title = slide_spec.get("title") or ""
    lede = slide_spec.get("lede") or slide_spec.get("body") or ""

    if eyebrow:
        _add_textbox(
            slide, x_px=SAFE_PAD, y_px=180, w_px=1500, h_px=24, text=eyebrow,
            font=FONT_SANS, size_pt=11, color=LC_RED, upper=True, letter_spacing_em=0.18,
        )
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=220, w_px=1600, h_px=320, text=title,
        font=FONT_SERIF, size_pt=64, color=LC_BLACK,
    )

    # 4px red rule + lede paragraph below
    _add_rect(slide, x_px=SAFE_PAD, y_px=620, w_px=4, h_px=180, fill=LC_RED)
    _add_textbox(
        slide, x_px=SAFE_PAD + 36, y_px=620, w_px=1500, h_px=300,
        text=lede, font=FONT_SANS, size_pt=22, color=LC_BLACK,
    )

    _page_no(slide, page_no)


# ── Custom: Structure (diagram PNG) ───────────────────────────────────


def _build_structure(prs: Presentation, slide_spec: dict, page_no: int):
    """Custom layout — eyebrow + diagram PNG centered on slide. The PNG is
    pre-rendered by frontend/scripts/render-diagram.js to a high-res image
    that fits the 1680×800 content area within the safe zone."""
    slide = _new_slide(prs)
    _chapter_tab(
        slide, num=f"{page_no:02d}",
        label="Recommended structure",
    )

    eyebrow = slide_spec.get("eyebrow") or "Recommended structure — entities and relationships"
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=140, w_px=1600, h_px=28, text=eyebrow,
        font=FONT_SANS, size_pt=11, color=LC_RED, upper=True, letter_spacing_em=0.18,
    )

    png_path = slide_spec.get("png_path")
    if png_path and Path(png_path).exists():
        # Fit the image into the canvas area below the eyebrow.
        # Available area: x=120, y=200; w=1680, h=820
        with Image.open(png_path) as img:
            img_w, img_h = img.size
        avail_w, avail_h = 1680, 820
        scale = min(avail_w / img_w, avail_h / img_h)
        draw_w = img_w * scale
        draw_h = img_h * scale
        # Centre within the available area
        draw_x = SAFE_PAD + (avail_w - draw_w) / 2
        draw_y = 200 + (avail_h - draw_h) / 2
        slide.shapes.add_picture(
            png_path, _px(draw_x), _px(draw_y), _px(draw_w), _px(draw_h),
        )
    else:
        _add_textbox(
            slide, x_px=SAFE_PAD, y_px=500, w_px=1680, h_px=80,
            text="(Structure diagram not available)",
            font=FONT_SANS, size_pt=18, color=LC_BLACK_60, align=PP_ALIGN.CENTER,
        )

    _page_no(slide, page_no)


# ── Layout 04: Disclaimer (verbatim) ──────────────────────────────────


_DISCLAIMER_LEFT = (
    "The contents of this document are confidential and are meant for the intended "
    "recipient only. If you are not the intended recipient, please delete all copies "
    "of this document and notify the sender immediately. This document, provided as "
    "a general commentary, is for informational purposes only and is not to be "
    "construed as an offer to sell or solicit an offer to buy any financial "
    "instruments in any jurisdiction. This does not constitute any form of "
    "regulated financial advice, and your independent financial advisor should be "
    "consulted prior to taking any investment decision(s).\n\n"
    "This document is based on information from sources which are reliable but has "
    "not been independently verified by Lighthouse Canton Pte. Ltd. (“LCPL”) and "
    "its affiliate, Lighthouse Canton Capital (DIFC) Pte. Ltd) (“LCC”) "
    "(collectively, “LC”). The Net Asset Value (“NAV”) of the funds will be "
    "computed independently by a third-party fund administrator on a regular basis "
    "in accordance with the applicable NAV calculation policy and relevant pricing "
    "model. Deviation to the policy, if any, will be disclosed by the respective "
    "fund administrator and/or the Manager. LC has taken the reasonable steps to "
    "verify the contents of this document and accepts no liability for any loss "
    "arising from the use of any information contained herein. Please also note "
    "that past performances are not indicative of future performance.\n\n"
    "Information contained herein are those of the author(s) and does not represent "
    "the views held by other parties. LC is also under no obligation to update you "
    "on any changes made to this document."
)

_DISCLAIMER_RIGHT = (
    "This document is prepared by LCPL and LCC, which are regulated by Monetary "
    "Authority of Singapore (“MAS”) and Dubai Financial Services Authority "
    "(“DFSA”) respectively. MAS and DFSA have no responsibility for reviewing, "
    "verifying and approving the contents of this document and/or other associated "
    "documents. The contents of this document may not be reproduced or referenced, "
    "either in part or in full, without prior written permission from LC.\n\n"
    "This document is only intended for Accredited Investors and/or Professional "
    "Clients, as defined by MAS and DFSA."
)

_DISCLAIMER_FOOTER = (
    "Regulated by MAS (Singapore), DFSA (Dubai), SEBI (India) and FCA (United Kingdom)."
)


def _build_disclaimer(prs: Presentation, page_no: int):
    """Layout 04 — disclaimer (verbatim from canonical-corporate-deck.html
    slide 03). MANDATORY for every external deck per SKILL.md §0."""
    slide = _new_slide(prs)
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=56, w_px=400, h_px=24, text="Disclaimer",
        font=FONT_SERIF, size_pt=14, color=LC_RED, italic=True,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=130, w_px=1600, h_px=80, text="Disclaimer",
        font=FONT_SERIF, size_pt=42, color=LC_BLACK,
    )

    col_w = (1920 - 2 * SAFE_PAD - 48) / 2  # 48px gap between columns
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=230, w_px=col_w, h_px=720,
        text=_DISCLAIMER_LEFT, font=FONT_SANS, size_pt=9, color=LC_BLACK,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD + col_w + 48, y_px=230, w_px=col_w, h_px=600,
        text=_DISCLAIMER_RIGHT, font=FONT_SANS, size_pt=9, color=LC_BLACK,
    )
    # Red rule + footer line on the right column
    _add_line(
        slide, x_px=SAFE_PAD + col_w + 48, y_px=900, w_px=col_w,
        color=LC_RED, weight_pt=1.5,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD + col_w + 48, y_px=910, w_px=col_w, h_px=30,
        text=_DISCLAIMER_FOOTER, font=FONT_SERIF, size_pt=10, color=LC_BLACK_80,
    )
    _page_no(slide, page_no)


# ── Layout 15: Offices (verbatim grid) ────────────────────────────────


_OFFICES = [
    ("Singapore", "HQ", "16 Collyer Quay #11-02\nCollyer Quay Centre\nSingapore–049318\n\n+65 67130570"),
    ("Dubai", "", "The Exchange,\nGate Village 11, Unit 204,\nDIFC, Dubai, UAE–507026\n\n+971 45 861500"),
    ("London", "", "24 Hanover Square,\nLondon, W1S 1JD\nUnited Kingdom"),
    ("New Delhi", "", "Unit 104A, Worldmark 2 Asset,\nDelhi Aerocity,\nNew Delhi 110037\n\n+91 9650473961"),
    ("Mumbai", "", "Unit 507/508, A Wing,\nINS Tower, G Block, BKC,\nMumbai 400051"),
    ("Bengaluru", "", "1st Floor, WeWork 37,\nCunningham Cross Rd,\nSRT Road, Vasant Nagar,\nBengaluru 560001\n\n+91 9900096873"),
    ("Hyderabad", "", "Suite 502, Building 450,\nCentral Plaza, Genome Valley,\nShameerpet, Hyderabad 500078\n\n+91 9900096873"),
    ("Chennai", "", "RK Swamy Centre, Hansa Building,\nDoor No:3, Thousand Lights,\nChennai 600006\n\n+91 9650473961"),
    ("GIFT City", "", "Unit FF-10, FF Floor,\nPragya Accelerator, Block 15T,\nGIFT CITY, Gandhinagar,\nGujarat 382355\n\n+91 9650473961"),
]

_GET_IN_TOUCH = "info@lighthouse-canton.com  ·  service@lighthouse-canton.in  ·  lighthouse-canton.com"


def _build_offices(prs: Presentation, page_no: int):
    """Layout 15 — offices grid + Get-in-touch row. MANDATORY closing slide."""
    slide = _new_slide(prs)
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=SAFE_PAD, w_px=400, h_px=24, text="Where we are",
        font=FONT_SERIF, size_pt=14, color=LC_RED, italic=True,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=160, w_px=1600, h_px=80, text="Offices",
        font=FONT_SERIF, size_pt=48, color=LC_BLACK,
    )

    # 4 columns × ceil(N/4) rows. With 9 cities that's 3 rows; the previous
    # touch_y math assumed 2 rows and crashed the Get-in-touch rule into
    # the third row's GIFT City cell.
    import math as _m
    grid_y = 270
    cell_w = (1920 - 2 * SAFE_PAD - 3 * 36) / 4
    row_h = 170          # tightened from 200 so 3 rows + closing block fit
    row_gap = 18
    n_rows = _m.ceil(len(_OFFICES) / 4)
    for i, (city, sub, body) in enumerate(_OFFICES):
        col = i % 4
        row = i // 4
        x = SAFE_PAD + col * (cell_w + 36)
        y = grid_y + row * (row_h + row_gap)
        _add_textbox(
            slide, x_px=x, y_px=y, w_px=cell_w, h_px=28,
            text=f"{city}{('  ' + sub) if sub else ''}",
            font=FONT_SERIF, size_pt=18, color=LC_RED, italic=True,
        )
        # Red underline rule
        _add_rect(slide, x_px=x, y_px=y + 30, w_px=cell_w * 0.4, h_px=2, fill=LC_RED)
        _add_textbox(
            slide, x_px=x, y_px=y + 42, w_px=cell_w, h_px=row_h - 48,
            text=body, font=FONT_SANS, size_pt=10, color=LC_BLACK,
        )

    # Get-in-touch row (full-width, red top border) — positioned BELOW the
    # actual last grid row, not a hard-coded y assuming 2 rows.
    touch_y = grid_y + n_rows * (row_h + row_gap) + 4
    _add_line(slide, x_px=SAFE_PAD, y_px=touch_y, w_px=1920 - 2 * SAFE_PAD, color=LC_RED, weight_pt=2)
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=touch_y + 14, w_px=300, h_px=30,
        text="Get in touch", font=FONT_SERIF, size_pt=18, color=LC_RED, italic=True,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD + 320, y_px=touch_y + 18, w_px=1500, h_px=30,
        text=_GET_IN_TOUCH, font=FONT_SERIF, size_pt=18, color=LC_BLACK,
    )

    _page_no(slide, page_no)


# ── Layout 13: Data table — used for client profile ───────────────────


def _build_table(prs: Presentation, slide_spec: dict, page_no: int):
    """Layout 13 — data table. We use this for the client profile (label/value
    rows). spec.rows = [{label, value}, ...]
    """
    slide = _new_slide(prs)
    _chapter_tab(
        slide, num=f"{page_no:02d}",
        label=slide_spec.get("chapter") or "Client",
    )
    eyebrow = slide_spec.get("eyebrow") or "Engagement"
    title = slide_spec.get("title") or "Client profile"

    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=140, w_px=1500, h_px=24, text=eyebrow,
        font=FONT_SANS, size_pt=11, color=LC_RED, upper=True, letter_spacing_em=0.18,
    )
    _add_textbox(
        slide, x_px=SAFE_PAD, y_px=176, w_px=1600, h_px=120, text=title,
        font=FONT_SERIF, size_pt=52, color=LC_BLACK,
    )

    rows = slide_spec.get("rows") or []
    row_y = 360
    row_h = 88
    for row in rows:
        label = row.get("label") or ""
        value = row.get("value") or row.get("body") or "—"
        _add_textbox(
            slide, x_px=SAFE_PAD, y_px=row_y + 12, w_px=300, h_px=row_h,
            text=label, font=FONT_SERIF, size_pt=22, color=LC_BLACK,
        )
        _add_rect(slide, x_px=SAFE_PAD + 320, y_px=row_y + 14, w_px=3, h_px=48, fill=LC_RED)
        _add_textbox(
            slide, x_px=SAFE_PAD + 360, y_px=row_y + 12, w_px=1300, h_px=row_h,
            text=value, font=FONT_SANS, size_pt=18, color=LC_BLACK_80,
        )
        # Bottom hairline
        _add_rect(slide, x_px=SAFE_PAD, y_px=row_y + row_h, w_px=1680, h_px=1, fill=LC_INK_400)
        row_y += row_h + 8

    _page_no(slide, page_no)


# ── Public API ────────────────────────────────────────────────────────


_BUILDERS = {
    "cover": _build_cover,
    "01": _build_cover,
    "02": _build_cover,
    "06": _build_pillars,
    "pillars": _build_pillars,
    "07": _build_title_lede,
    "title-lede": _build_title_lede,
    "13": _build_table,
    "table": _build_table,
    "profile": _build_table,
    "structure": _build_structure,
}


def build_pptx(spec: dict, out_path: Path) -> Path:
    """Build a .pptx from a curator JSON spec.

    spec = {
      "client_name": str,
      "slides": [
        {"layout": "cover" | "06" | "07" | "structure" | "04" | "15" | ...,
         ...layout-specific fields...},
        ...
      ]
    }

    Disclaimer (04) and Offices (15) are appended automatically if not
    already in the spec — they are mandatory per SKILL.md §0.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W_EMU
    prs.slide_height = SLIDE_H_EMU

    slides = list(spec.get("slides") or [])
    layouts_present = {s.get("layout") for s in slides}
    if "04" not in layouts_present and "disclaimer" not in layouts_present:
        slides.append({"layout": "04"})
    if "15" not in layouts_present and "offices" not in layouts_present:
        slides.append({"layout": "15"})

    page_no = 1
    for slide_spec in slides:
        layout = (slide_spec.get("layout") or "").lower()
        if layout in ("04", "disclaimer"):
            _build_disclaimer(prs, page_no)
        elif layout in ("15", "offices"):
            _build_offices(prs, page_no)
        elif layout in _BUILDERS:
            _BUILDERS[layout](prs, slide_spec, page_no)
        else:
            logger.warning("Unknown layout %r — falling back to title-lede", layout)
            _build_title_lede(prs, slide_spec, page_no)
        page_no += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
